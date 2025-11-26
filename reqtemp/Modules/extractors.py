from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import re
import streamlit as st

def extract_block(tag, text):
    pattern = rf"<{tag}>(.*?)(?=<[A-Z][A-Z_ ]+>|$)"

    match = re.search(pattern, text, flags=re.DOTALL)
    if not match:
        return ""
    return match.group(1).strip()

def extract_text_from_file(uploaded_file):
    """Extract text from PDF, DOCX, XLSX, or PPTX."""
    text = ""

    # PDF
    if uploaded_file.name.lower().endswith(".pdf"):
        pdf = PdfReader(uploaded_file)
        for page in pdf.pages:
            try:
                text += page.extract_text() + "\n"
            except:
                continue

    # DOCX
    elif uploaded_file.name.lower().endswith(".docx"):
        doc = Document(uploaded_file)
        text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())

    # XLSX → convert entire sheet to comma-separated text
    elif uploaded_file.name.lower().endswith(".xlsx"):
        import pandas as pd
        df = pd.read_excel(uploaded_file).fillna("")
        text = df.to_csv(index=False)

    # PPTX → extract text from all slides
    elif uploaded_file.name.lower().endswith(".pptx"):
        from pptx import Presentation
        pres = Presentation(uploaded_file)
        slides_text = []

        for slide in pres.slides:
            for shape in slide.shapes:

                # Normal text
                if hasattr(shape, "text") and shape.text.strip():
                    slides_text.append(shape.text.strip())

                # Table text
                if shape.has_table:
                    table = shape.table
                    tbl_rows = []
                    for row in table.rows:
                        cells = [c.text.strip() for c in row.cells]
                        tbl_rows.append(" | ".join(cells))
                    slides_text.append("\n".join(tbl_rows))

        text = "\n".join(slides_text)

    else:
        st.warning("⚠️ Unsupported file type.")
        return ""

    return text.strip()

from pptx import Presentation

def extract_image_from_slide(pptx_file, slide_number):
    prs = Presentation(pptx_file)

    # PPT slide index is 0-based → slide 17 means index 16
    if slide_number - 1 >= len(prs.slides):
        return None

    slide = prs.slides[slide_number - 1]

    # Scan all shapes in the slide
    for shape in slide.shapes:
        if shape.shape_type == 13:  # 13 = PICTURE
            image = shape.image
            return image.blob, image.ext  # binary content + file extension

    return None

def extract_table_from_slide(uploaded_pptx, slide_number):
    """
    Extracts the FIRST table found in the given slide.
    Returns a markdown-like table string.
    """

    from pptx import Presentation

    prs = Presentation(uploaded_pptx)

    if slide_number < 1 or slide_number > len(prs.slides):
        return None

    slide = prs.slides[slide_number - 1]

    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table

            # Extract text from table → convert to markdown-style table
            rows = []
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                rows.append("| " + " | ".join(cells) + " |")

            # Insert header separator for markdown
            if len(rows) > 1:
                header = rows[0]
                col_count = header.count("|") - 1
                separator = "| " + " | ".join(["---"] * col_count) + " |"
                rows.insert(1, separator)

            return "\n".join(rows)

    return None

def extract_slide7_summary(uploaded_pptx):
    """
    Extract ONLY bullet points from slide 7.
    Supports weird PPT line breaks (vertical tab: \x0b).
    """
    from pptx import Presentation
    prs = Presentation(uploaded_pptx)

    slide_number = 7
    if slide_number < 1 or slide_number > len(prs.slides):
        return ""

    slide = prs.slides[slide_number - 1]
    bullets = []

    for shape in slide.shapes:

        # Skip tables
        if shape.has_table:
            continue

        if hasattr(shape, "text") and shape.text.strip():
            raw = shape.text

            # 🔥 Replace PPT weird line-break with normal newline
            raw = raw.replace("\x0b", "\n")

            # Split into lines
            lines = raw.split("\n")

            for line in lines:
                stripped = line.strip()

                # Detect bullets that begin with -, •, – OR have "- " later in the line
                if (
                    stripped.startswith("-")
                    or stripped.startswith("•")
                    or stripped.startswith("–")
                    or stripped.startswith("—")
                    or stripped.startswith("·")
                ):
                    bullets.append(stripped)

                # Also catch text like:
                # "Based on our migration assessment:\n- 71.49% ..."
                elif "- " in stripped:
                    idx = stripped.find("- ")
                    bullets.append(stripped[idx:].strip())

    return "\n".join(bullets).strip()

import re

def extract_total_interfaces_from_slide(uploaded_pptx, slide_number=5):
    """
    Robust extraction of the total interfaces count from a slide.
    Strategy:
      - collect all text (text boxes + table cells)
      - debug-print a snippet so you can inspect what's been captured
      - first, try targeted regex patterns near phrases like 'No of Interfaces'
      - fallback: find all numbers on the slide, pick the largest (ignoring tiny serials where possible)
    Returns a string number (no commas) or None.
    """

    from pptx import Presentation

    prs = Presentation(uploaded_pptx)

    if slide_number < 1 or slide_number > len(prs.slides):
        return None

    slide = prs.slides[slide_number - 1]

    all_text = []

    for shape in slide.shapes:
        # TEXT
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            txt = shape.text.strip()
            # replace weird vertical/tab breaks
            txt = txt.replace("\x0b", " ").replace("\n", " ")
            all_text.append(txt)

        # TABLE
        if shape.has_table:
            tbl = shape.table
            for row in tbl.rows:
                for cell in row.cells:
                    if cell.text and cell.text.strip():
                        cell_txt = cell.text.strip().replace("\n", " ")
                        all_text.append(cell_txt)

    combined = " ".join(all_text).strip()

    # Try targeted patterns first (handles commas and spaces)
    patterns = [
        r"(?:No\.?\s*of\s*Interfaces(?:\s*in\s*Production|\s*to\s*be\s*migrated)?)[\:\-\s]*([0-9][0-9,]*)",
        r"(?:No\.?\s*of\s*Interfaces)[\:\-\s]*([0-9][0-9,]*)",
        r"(?:Total\s+Number\s+of\s+Interfaces)[\:\-\s]*([0-9][0-9,]*)",
        r"(?:No\s+of\s+Interfaces\s+to\s+be\s+migrated)[\:\-\s]*([0-9][0-9,]*)",
        r"\bTotal\s*[:\-]?\s*([0-9][0-9,]*)\b"
    ]

    for pat in patterns:
        m = re.search(pat, combined, flags=re.I)
        if m:
            val = m.group(1)
            val = val.replace(",", "")
            return val

    # Fallback: find all numeric tokens, clean them and pick the largest sensible one
    nums = re.findall(r"\d[\d,]*", combined)
    cleaned = []
    for n in nums:
        n_clean = n.replace(",", "")
        try:
            cleaned.append(int(n_clean))
        except:
            continue

    if not cleaned:
        return None

    # Heuristic: prefer numbers > 10 (to avoid serial numbers). If none > 10 exist, choose the max.
    large_candidates = [v for v in cleaned if v > 10]
    if large_candidates:
        return str(max(large_candidates))
    else:
        return str(max(cleaned))

def detect_client_name_from_text(text: str) -> str:
    """
    Smarter detection of client name from RFP or SOW text.
    Scans for patterns like 'Client:', 'Prepared for', 'RFP from', 'Issued by', etc.
    Falls back gracefully if not found.
    """
    patterns = [
        r"(?i)\bclient\s*(?:name)?\s*[:\-]\s*([A-Za-z0-9&,\.\s]+)",
        r"(?i)\bprepared\s*for\s*([A-Za-z0-9&,\.\s]+)",
        r"(?i)\bproposal\s*(?:for|to)\s*([A-Za-z0-9&,\.\s]+)",
        r"(?i)\brfp\s*(?:from|by|for)\s*([A-Za-z0-9&,\.\s]+)",
        r"(?i)\bissued\s*(?:by|to)\s*([A-Za-z0-9&,\.\s]+)",
        r"(?i)\bsubmitted\s*(?:by|to)\s*([A-Za-z0-9&,\.\s]+)",
        r"(?i)\borganization\s*[:\-]\s*([A-Za-z0-9&,\.\s]+)",
    ]

    for pat in patterns:
        match = re.search(pat, text)
        if match:
            name = match.group(1).strip()
            # clean any trailing words like "Limited", "LLC", etc.
            name = re.sub(r"\s+(Limited|Ltd|LLC|Company|Inc\.?)\b.*", r" \1", name, flags=re.I)
            # remove extra newlines or dots
            name = re.sub(r"[\n\r]+", " ", name).strip(" .")
            # cap first letters
            return name.title()

    return "Client"


def summarize_large_rfp(client, model_name, text, max_chunk_size=3500):
    """
    Automatically handles large RFPs by splitting into chunks,
    summarizing each chunk, and returning a final merged summary.
    """

    # Split into ~3500-character chunks
    chunks = []
    current = ""

    for line in text.splitlines():
        if len(current) + len(line) > max_chunk_size:
            chunks.append(current)
            current = ""
        current += line + "\n"

    if current:
        chunks.append(current)

    st.info(f"🔍 RFP too large — splitting into {len(chunks)} chunks for safe processing.")

    summaries = []
    progress = st.progress(0)
    status = st.empty()

    total_chunks = len(chunks)

    for i, chunk in enumerate(chunks, start=1):

        # Show progress message
        status.write(f"⏳ Summarizing chunk {i}/{total_chunks}…")

        # Update progress bar (value between 0 and 1)
        progress.progress(i / total_chunks)

        prompt = f"""
You are an SAP domain expert. Summarize the following RFP section clearly and concisely.
Focus only on project requirements, scope, process details, integrations, and constraints.

CHUNK {i}:
{chunk}
"""

        response = client.chat.completions.create(
            model=model_name,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.2
        )

        summaries.append(response.choices[0].message.content.strip())

    # Combine all part summaries into one master text
    final_summary = "\n\n".join(summaries)
    progress.progress(1.0)
    st.success("📘 RFP summarized into a clean master document.")

    return final_summary

